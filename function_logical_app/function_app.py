import json
import logging
import re
import base64, io, os, textwrap
from typing import List
import azure.functions as func
from dotenv import load_dotenv
import os
import time
from azure.identity import DefaultAzureCredential
from azure.ai.projects import AIProjectClient
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import ast
# Deployment : Jun 4 : 9.35 AM


from openpyxl import load_workbook, Workbook

load_dotenv()

app = func.FunctionApp(http_auth_level=func.AuthLevel.ANONYMOUS)


@app.route(route="health_check")
def health_check(req: func.HttpRequest) -> func.HttpResponse:
    return func.HttpResponse(
        json.dumps({"content": "ok"}), status_code=200, mimetype="application/json"
    )


def call_agent(prompt: str) -> dict:
    endpoint = "https://taqa-electric-grid.services.ai.azure.com/api/projects/taqa-electric-grid"
    agent_id = "asst_NbiJUw9yQGRWX1LLMC7dWBdI"
    AZURE_FOUNDRY_AGENTS_APIKEY = "8YHvIOstZSLqX7Lb8plm4k2XqLRWQSDWGGJ4ypuQDpZqKyWM74H3JQQJ99BEACHYHv6XJ3w3AAAAACOGYX4y"
    if not endpoint or not agent_id:
        raise RuntimeError(
            "Environment variables AZURE_AI_PROJECT_ENDPOINT and AGENT_ID must be set."
        )

    client = AIProjectClient(
        endpoint=endpoint,
        credential=DefaultAzureCredential(),
        # api_version="latest"
    )

    thread = client.agents.threads.create()
    client.agents.messages.create(thread_id=thread.id, role="user", content=prompt)
    run = client.agents.runs.create_and_process(thread_id=thread.id, agent_id=agent_id)
    while run.status in ("queued", "in_progress"):
        time.sleep(1)
        run = client.agents.runs.get(thread_id=thread.id, run_id=run.id)

    all_msgs = list(client.agents.messages.list(thread_id=thread.id))

    assistant_msgs = [m for m in all_msgs if m.role == "assistant"]

    if not assistant_msgs:
        raise RuntimeError("No assistant message found in thread.")
    last_assistant = assistant_msgs[-1]

    content_obj = last_assistant.content
    if hasattr(content_obj, "text"):
        assistant_reply = content_obj.text
    else:
        assistant_reply = str(content_obj)
    return {"response": assistant_reply}


def parse_progress_from_response(response_str):
    match = re.search(r'"progress"\s*:\s*([0-9]+(?:\.[0-9]+)?)', response_str)
    if not match:
        raise ValueError("Progress value not found in the response string.")
    progress_value = float(match.group(1))
    return {"progress": progress_value}


def load_workbook_from_base64(b64_string: str) -> Workbook:
    """
    Decode a Base64 string into bytes, then load it into an openpyxl Workbook.
    Returns the loaded Workbook object (preserving all formatting).
    """
    file_bytes = base64.b64decode(b64_string)
    stream = io.BytesIO(file_bytes)
    wb = load_workbook(filename=stream, data_only=False)
    return wb


def save_workbook_to_base64(wb: Workbook) -> str:
    stream = io.BytesIO()
    wb.save(stream)
    stream.seek(0)
    return base64.b64encode(stream.read()).decode("utf-8")


@app.route(route="prcoess_excel_file")
@app.function_name(name="prcoess_excel_file")
def prcoess_excel_file(req: func.HttpRequest) -> func.HttpResponse:
    req_body = req.get_json()
    content_b64 = req_body.get("content")
    # file_bytes = base64.b64decode(content_b64)
    # output_path = os.path.join(os.getcwd(), "output.xlsx")

    # with open(output_path, "wb") as f:
    #         f.write(file_bytes)

    user_prompt = """
                    You are given access to a tool named `extractdatafromwordfile_Tool`. When instructed, you must:  
                    1. Call `extractdatafromwordfile_Tool` to retrieve its output (a JSON object containing a field called `"progress"` whose value is in percentage  
                    2. Read the `"From Over all progress planned you have to find out actual "` Find out this form the output of the tool.  
                    3. Return only a JSON object with a single key `"progress"` and its   
                    value set to the computed percentage (as a number). Do not include any other keys or any extra text.  

                    For example, if the tool returns:  
                    ```json  
                    {  
                    "progress": "the overall progress actual",  
                    }  

                  """
    result = call_agent(user_prompt)
    progress = parse_progress_from_response(result["response"])
    wb = load_workbook_from_base64(content_b64)

    # 4) Access the "Progress %" sheet
    sheet_name = "Progress %"
    if sheet_name not in wb.sheetnames:
        return func.HttpResponse(
            json.dumps({"error": f"Worksheet '{sheet_name}' not found."}),
            status_code=400,
            mimetype="application/json",
        )
    ws = wb[sheet_name]

    ws["B4"].value = 1
    ws["C4"].value = "Project 1"
    ws["D4"].value = progress["progress"]
    # wb.save("input.xlsx")

    new_b64 = save_workbook_to_base64(wb)

    # Decode the Base64 string into binary Excel content
    excel_data = base64.b64decode(new_b64)

    # Write the binary content to an Excel file
    with open("output.xlsx", "wb") as f:
        f.write(excel_data)
    return func.HttpResponse(
        json.dumps({"content": new_b64}), status_code=200, mimetype="application/json"
    )


def parse_date(date_str):
    try:
        return datetime.strptime(date_str.strip(), "%d.%m.%Y").date()
    except:
        return None


def update_slide2_table(table, expected_value):
    for row in table.rows:
        if len(row.cells) < 2:
            continue

        label = row.cells[0].text.strip().lower()

        # if "awarded value" in label:
        #     row.cells[1].text = awarded_value
        #     for paragraph in row.cells[1].text_frame.paragraphs:
        #         for run in paragraph.runs:
        #             run.font.name = "Arial"
        #             run.font.size = Pt(16)

        # elif "expected completion date" in label:
        #     row.cells[1].text = expected_value
        #     for paragraph in row.cells[1].text_frame.paragraphs:
        #         for run in paragraph.runs:
        #             run.font.name = "Arial"
        #             run.font.size = Pt(16)

        if "expected completion date" in label:
            row.cells[1].text = expected_value
            for paragraph in row.cells[1].text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(16)


def update_slide3_table(table, date_lists):
    expected_actual_cols = [2, 4, 6, 8, 10, 12]
    plan_cols = [1, 3, 5, 7, 9, 11]

    for row_offset in range(5):
        row = table.rows[row_offset + 2]
        for i in range(6):
            col_idx = expected_actual_cols[i]
            plan_idx = plan_cols[i]

            actual_date_str = date_lists[i][row_offset].strip()
            actual_date = parse_date(actual_date_str)
            plan_text = row.cells[plan_idx].text.strip()
            plan_date = parse_date(plan_text)

            cell = row.cells[col_idx]
            cell.text = ""
            paragraph = cell.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = actual_date_str

            font = run.font
            font.name = "Arial"
            font.size = Pt(11)

            if actual_date and plan_date:
                if actual_date < plan_date:
                    font.color.rgb = RGBColor(0, 128, 0)  # Green
                elif actual_date == plan_date:
                    font.color.rgb = RGBColor(0, 0, 255)  # Blue
                else:
                    font.color.rgb = RGBColor(255, 0, 0)  # Red
            else:
                font.color.rgb = RGBColor(0, 0, 0)


def process_presentation_base64(content_b64, expected_value, date_lists):
    ppt_bytes = base64.b64decode(content_b64)
    prs = Presentation(io.BytesIO(ppt_bytes))

    if len(prs.slides) > 2:
        # Slide 2 updates
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if shape.has_table:
                update_slide2_table(shape.table, expected_value)

        # Slide 3 updates
        slide3 = prs.slides[2]
        for shape in slide3.shapes:
            if shape.has_table:
                update_slide3_table(shape.table, date_lists)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    # filename = f"updated_ppt_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    # with open(filename, "wb") as f:
    #     f.write(output.getvalue())

    return {
        "status": "success",
        # "filename": filename,
        "content_b64": base64.b64encode(output.getvalue()).decode("utf-8"),
    }


def parse_completion_date_from_response(response_str):
    match = re.search(r"'value'\s*:\s*'([^']+)'", response_str)
    if not match:
        raise ValueError("Completion date not found in the response string.")
    return match.group(1)


Expected_completion_Date = """
                   Use the tool to retrieve the content. From the content,
                   extract the value corresponding to "Expected Completion Date".
                   Return only the date in `YYYY-MM-DD` format. Do not include any explanation, 
                   label, or extra words in the output. Just output the extracted date.
                  """


Date_Extraction_Prompt = """
                Call the tool `extractdatafromwordfile_Tool` to extract data from the Word file.
                The tool returns data in Markdown format as JSON. From this data, locate the table 
                that contains activity-related information, specifically one with columns for PLAN and Expected / Actual dates.

                From that table, extract only the Expected / Actual date values. Format each column
                as a list of 5 dates in the format DD.MM.YYYY. Return the final result as a list of 6 lists,
                where each list represents one column and contains exactly 5 date entries.
                 
                **Important**: Your final output must be in the format month values should be int if may convert it into 5 . 
                and split date with dot  like this 17.11.2023:
                [
                ["17.11.2023", "14.10.2024", "08.04.2025", "24.06.2025", "15.09.2025"],
                ["..."],
                ...
                ]
                **Only return the list of lists. Do not include any explanation, headers, labels, or other content.**
"""




@app.route(route="update_powerpoint")
@app.function_name(name="update_powerpoint")
def update_powerpoint(req: func.HttpRequest) -> func.HttpResponse:
    try:
        req_body = req.get_json()
        content_b64 = req_body.get("content")

        expected_value = call_agent(Expected_completion_Date)
        parsed_response = ast.literal_eval(expected_value["response"])
        completion_date = parsed_response[0]["text"]["value"]

        date_list = call_agent(Date_Extraction_Prompt)
        parsed_date = ast.literal_eval(date_list["response"])
        date_lists = ast.literal_eval(parsed_date[0]["text"]["value"])

        # expected_value=parse_completion_date_from_response(expected_value)
        # awarded_value = req_body.get("awarded_value", "AED 23000")

        # date_lists = req_body.get(
        #     "date_lists",
        #     [
        #         ["17.11.2023", "14.10.2024", "08.04.2025", "24.06.2025", "15.09.2025"],
        #         ["31.01.2024", "06.06.2024", "15.11.2024", "30.06.2025", "15.09.2025"],
        #         ["15.12.2023", "25.06.2024", "20.11.2024", "15.04.2025", "03.07.2025"],
        #         ["24.05.2024", "07.04.2025", "07.06.2025", "26.06.2025", "30.09.2025"],
        #         ["24.05.2024", "20.01.2025", "11.03.2025", "02.06.2025", "16.06.2025"],
        #         ["15.04.2024", "14.04.2025", "30.04.2025", "27.05.2025", "07.02.2025"],
        #     ],
        # )
        if len(date_lists) != 6 or any(len(sublist) != 5 for sublist in date_lists):
            raise ValueError("date_lists must contain 6 lists with 5 dates each")

        result = process_presentation_base64(content_b64, completion_date, date_lists)

        return func.HttpResponse(
            json.dumps(
                {
                    # "filename": result["filename"],
                    "content": result["content_b64"],
                    # "expected_value":completion_date,
                    # "date":date_lists
                }
            ),
            status_code=200,
            mimetype="application/json",
        )

    except Exception as e:
        return func.HttpResponse(
            json.dumps({"status": "error", "message": str(e)}),
            status_code=500,
            mimetype="application/json",
        )
