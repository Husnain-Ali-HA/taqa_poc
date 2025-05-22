# from odf.opendocument import load
# from odf.style import (
#     Style,
#     TableColumnProperties,
#     TableCellProperties,
#     ParagraphProperties,
#     TextProperties,
#     TableRowProperties,
# )
# from odf.table import Table, TableColumn, TableRow, TableCell
# from odf.text import P
# from odf.draw import Page, Frame


# def add_data_slide1(doc, summary: dict):
#     slides = doc.presentation.getElementsByType(Page)
#     if not slides:
#         raise RuntimeError("Presentation has no slides")
#     slide1 = slides[0]
#     styles = doc.automaticstyles

#     cell_style = "Slide1Cell"
#     if not any(
#         s.getAttribute("name") == cell_style for s in styles.getElementsByType(Style)
#     ):
#         st = Style(name=cell_style, family="table-cell")
#         st.addElement(TableCellProperties(backgroundcolor="#ffffff"))
#         st.addElement(ParagraphProperties(textalign="center"))
#         st.addElement(
#             TextProperties(fontfamily="Arial", fontsize="12pt", color="#000000")
#         )
#         styles.addElement(st)

#     col_style = "Slide1Col"
#     if not any(
#         s.getAttribute("name") == col_style for s in styles.getElementsByType(Style)
#     ):
#         stc = Style(name=col_style, family="table-column")
#         stc.addElement(TableColumnProperties(columnwidth="5cm"))
#         styles.addElement(stc)

#     # — build table —
#     table = Table(name="SummaryTable")

#     table.addElement(TableColumn(stylename=col_style))
#     table.addElement(TableColumn(stylename=col_style))

#     for key, value in summary.items():
#         tr = TableRow()
#         for text in (key, value):
#             tc = TableCell(stylename=cell_style)
#             tc.addElement(P(text=str(text)))
#             tr.addElement(tc)
#         table.addElement(tr)

#     frame = Frame(
#         anchortype="page",
#         x="8cm",
#         y="6cm",
#         width="10cm",
#         height=f"{len(summary) * 0.8 + 0.5}cm",
#     )
#     frame.addElement(table)
#     slide1.addElement(frame)


# def modify_slide2(doc, data):
#     SLIDE_WIDTH_CM = 28.0
#     SLIDE_HEIGHT_CM = 21.0
#     MARGIN_LEFT_CM = 0.0
#     MARGIN_TOP_CM = 0.0
#     FRAME_WIDTH_CM = SLIDE_WIDTH_CM - 2 * MARGIN_LEFT_CM
#     FRAME_HEIGHT_CM = SLIDE_HEIGHT_CM - 2 * MARGIN_TOP_CM

#     num_rows = len(data)
#     row_height_str = f"{FRAME_HEIGHT_CM / num_rows:.2f}cm"

#     COL_PROPS = [3, 6, 2, 3, 3]
#     total_prop = sum(COL_PROPS)
#     col_widths = [f"{FRAME_WIDTH_CM * prop / total_prop:.2f}cm" for prop in COL_PROPS]

#     styles = doc.automaticstyles

#     for idx, width in enumerate(col_widths):
#         name = f"TableCol{idx}"
#         if not any(
#             s.getAttribute("name") == name for s in styles.getElementsByType(Style)
#         ):
#             st = Style(name=name, family="table-column")
#             st.addElement(TableColumnProperties(columnwidth=width))
#             styles.addElement(st)

#     hdr_name = "CustomHeader"
#     if not any(
#         s.getAttribute("name") == hdr_name for s in styles.getElementsByType(Style)
#     ):
#         hdr = Style(name=hdr_name, family="table-cell")
#         hdr.addElement(TableCellProperties(backgroundcolor="#4f81bd", padding="0.1cm"))
#         hdr.addElement(ParagraphProperties(textalign="center"))
#         hdr.addElement(
#             TextProperties(
#                 fontfamily="Arial", fontweight="bold", color="#ffffff", fontsize="12pt"
#             )
#         )
#         styles.addElement(hdr)

#     default_name = "DefaultCell"
#     if not any(
#         s.getAttribute("name") == default_name for s in styles.getElementsByType(Style)
#     ):
#         dft = Style(name=default_name, family="table-cell")
#         dft.addElement(ParagraphProperties(textalign="left", marginleft="0.2cm"))
#         dft.addElement(
#             TextProperties(fontfamily="Arial", fontsize="10pt", color="#000000")
#         )
#         styles.addElement(dft)

#     num_name = "CustomNumeric"
#     if not any(
#         s.getAttribute("name") == num_name for s in styles.getElementsByType(Style)
#     ):
#         num = Style(name=num_name, family="table-cell")
#         num.addElement(ParagraphProperties(textalign="right", marginright="0.2cm"))
#         num.addElement(
#             TextProperties(fontfamily="Arial", fontsize="10pt", color="#000000")
#         )
#         styles.addElement(num)

#     row_style = "DynamicRow"
#     if not any(
#         s.getAttribute("name") == row_style for s in styles.getElementsByType(Style)
#     ):
#         rs = Style(name=row_style, family="table-row")
#         rs.addElement(
#             TableRowProperties(rowheight=row_height_str, minrowheight=row_height_str)
#         )
#         styles.addElement(rs)

#     slides = doc.presentation.getElementsByType(Page)
#     if len(slides) < 2:
#         raise RuntimeError("Presentation must have at least 2 slides")
#     slide2 = slides[1]

#     table = Table(name="LineItemTable")
#     for idx in range(len(col_widths)):
#         table.addElement(TableColumn(stylename=f"TableCol{idx}"))

#     for row_idx, row in enumerate(data):
#         tr = TableRow(stylename=row_style)
#         for col_idx, text in enumerate(row):
#             tc = TableCell()
#             if row_idx == 0:
#                 tc.setAttribute("stylename", hdr_name)
#             elif col_idx in (2, 3, 4):
#                 tc.setAttribute("stylename", num_name)
#             else:
#                 tc.setAttribute("stylename", default_name)
#             tc.addElement(P(text=str(text)))
#             tr.addElement(tc)
#         table.addElement(tr)

#     frame = Frame(
#         anchortype="page",
#         x=f"{MARGIN_LEFT_CM}cm",
#         y=f"{MARGIN_TOP_CM}cm",
#         width=f"{FRAME_WIDTH_CM}cm",
#         height=f"{FRAME_HEIGHT_CM}cm",
#     )
#     frame.addElement(table)
#     slide2.addElement(frame)


# def modify_existing_odp(
#     input_path: str, output_path: str, table_data: list, summary_json: dict
# ):
#     doc = load(input_path)
#     add_data_slide1(doc, summary_json)
#     modify_slide2(doc, table_data)
#     doc.save(output_path)

# from odf.opendocument import load
# from odf.style import (
#     Style,
#     TableColumnProperties,
#     TableCellProperties,
#     ParagraphProperties,
#     TextProperties,
#     TableRowProperties,
# )
# from odf.table import Table, TableColumn, TableRow, TableCell
# from odf.text import P
# from odf.draw import Page, Frame


# def add_data_slide1(doc, summary: dict):
#     slides = doc.presentation.getElementsByType(Page)
#     if not slides:
#         raise RuntimeError("Presentation has no slides")
#     slide1 = slides[0]
#     styles = doc.automaticstyles

#     # white-cell style, centered text
#     cell_style = "Slide1Cell"
#     if not any(s.getAttribute("name")==cell_style for s in styles.getElementsByType(Style)):
#         st = Style(name=cell_style, family="table-cell")
#         st.addElement(TableCellProperties(backgroundcolor="#ffffff", verticalalign="middle"))
#         st.addElement(ParagraphProperties(textalign="center"))
#         st.addElement(TextProperties(fontfamily="Arial", fontsize="12pt", color="#000000"))
#         styles.addElement(st)

#     # two equal columns
#     col_style = "Slide1Col"
#     if not any(s.getAttribute("name")==col_style for s in styles.getElementsByType(Style)):
#         stc = Style(name=col_style, family="table-column")
#         stc.addElement(TableColumnProperties(columnwidth="5cm"))
#         styles.addElement(stc)

#     # build 2-column table
#     table = Table(name="SummaryTable")
#     table.addElement(TableColumn(stylename=col_style))
#     table.addElement(TableColumn(stylename=col_style))

#     for key, value in summary.items():
#         tr = TableRow()
#         for text in (key, value):
#             tc = TableCell(stylename=cell_style)
#             tc.addElement(P(text=str(text)))
#             tr.addElement(tc)
#         table.addElement(tr)

#     # center on slide
#     SLIDE_W, SLIDE_H = 28.0, 21.0  # cm
#     rows = len(summary)
#     frame_w, frame_h = 5.0*2, rows*0.8 + 0.5
#     x = (SLIDE_W - frame_w) / 2
#     y = (SLIDE_H - frame_h) / 2

#     frame = Frame(
#         anchortype="page",
#         x=f"{x:.2f}cm",
#         y=f"{y:.2f}cm",
#         width=f"{frame_w:.2f}cm",
#         height=f"{frame_h:.2f}cm"
#     )
#     frame.addElement(table)
#     slide1.addElement(frame)


# def modify_slide2(doc, data):
#     SLIDE_W, SLIDE_H = 25.0, 21.0
#     FRAME_W = 25
#     FRAME_H = 21

#     rows = len(data)
#     row_h = FRAME_H/rows
#     row_h_str = f"{row_h:.2f}cm"

#     COL_PROPS = [3,6,2,3,3]
#     total = sum(COL_PROPS)
#     col_widths = [f"{FRAME_W*prop/total:.2f}cm" for prop in COL_PROPS]

#     styles = doc.automaticstyles

#     # column styles
#     for i, w in enumerate(col_widths):
#         name = f"TableCol{i}"
#         if not any(s.getAttribute("name")==name for s in styles.getElementsByType(Style)):
#             st = Style(name=name, family="table-column")
#             st.addElement(TableColumnProperties(columnwidth=w))
#             styles.addElement(st)

#     # header style
#     hdr = "CustomHeader"
#     if not any(s.getAttribute("name")==hdr for s in styles.getElementsByType(Style)):
#         h = Style(name=hdr, family="table-cell")
#         h.addElement(TableCellProperties(backgroundcolor="#4f81bd", padding="0.1cm"))
#         h.addElement(ParagraphProperties(textalign="center"))
#         h.addElement(TextProperties(fontfamily="Arial", fontweight="bold",
#                                    color="#ffffff", fontsize="12pt"))
#         styles.addElement(h)

#     # default style
#     dft = "DefaultCell"
#     if not any(s.getAttribute("name")==dft for s in styles.getElementsByType(Style)):
#         d = Style(name=dft, family="table-cell")
#         d.addElement(ParagraphProperties(textalign="left", marginleft="0.2cm"))
#         d.addElement(TextProperties(fontfamily="Arial", fontsize="10pt", color="#000000"))
#         styles.addElement(d)

#     # numeric style
#     num = "CustomNumeric"
#     if not any(s.getAttribute("name")==num for s in styles.getElementsByType(Style)):
#         n = Style(name=num, family="table-cell")
#         n.addElement(ParagraphProperties(textalign="right", marginright="0.2cm"))
#         n.addElement(TextProperties(fontfamily="Arial", fontsize="10pt", color="#000000"))
#         styles.addElement(n)

#     # row style
#     rsty = "DynamicRow"
#     if not any(s.getAttribute("name")==rsty for s in styles.getElementsByType(Style)):
#         rs = Style(name=rsty, family="table-row")
#         rs.addElement(TableRowProperties(rowheight=row_h_str, minrowheight=row_h_str))
#         styles.addElement(rs)

#     slides = doc.presentation.getElementsByType(Page)
#     if len(slides)<2:
#         raise RuntimeError("Need at least 2 slides")
#     slide2 = slides[1]

#     table = Table(name="LineItemTable")
#     for i in range(len(col_widths)):
#         table.addElement(TableColumn(stylename=f"TableCol{i}"))

#     for ridx, row in enumerate(data):
#         tr = TableRow(stylename=rsty)
#         for cidx, txt in enumerate(row):
#             tc = TableCell()
#             if ridx==0:
#                 tc.setAttribute("stylename", hdr)
#             elif cidx in (2,3,4):
#                 tc.setAttribute("stylename", num)
#             else:
#                 tc.setAttribute("stylename", dft)
#             tc.addElement(P(text=str(txt)))
#             tr.addElement(tc)
#         table.addElement(tr)

#     frame = Frame(
#         anchortype="page",
#         x="0cm",
#         y="0cm",
#         width=f"{FRAME_W}cm",
#         height=f"{FRAME_H}cm"
#     )
#     frame.addElement(table)
#     slide2.addElement(frame)


# def modify_existing_odp(input_path: str,
#                         table_data: list,
#                         summary_json: dict,
#                         output_path: str):
#     doc = load(input_path)
#     add_data_slide1(doc, summary_json)
#     modify_slide2(doc, table_data)

#     if not isinstance(output_path, str) or not output_path:
#         raise ValueError("Invalid output_path")
#     doc.save(output_path)
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_data_slide1(prs: Presentation, summary: dict):
    slide = prs.slides[0]
    rows, cols = len(summary), 2

    col_w, frame_h = 5.0, rows * 0.8 + 0.5
    frame_w = col_w * cols

    left = (prs.slide_width - Cm(frame_w)) / 2
    top  = (prs.slide_height - Cm(frame_h)) / 2

    tbl = slide.shapes.add_table(rows, cols, left, top, Cm(frame_w), Cm(frame_h)).table
    for c in tbl.columns:
        c.width = Cm(col_w)

    for r, (k, v) in enumerate(summary.items()):
        for c, txt in enumerate((k, v)):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p = cell.text_frame.paragraphs[0]
            p.text = str(txt)
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

def modify_slide2(prs: Presentation, data: list):
    slide = prs.slides[1]
    rows, cols = len(data), len(data[0])

    slide_w, slide_h = prs.slide_width, prs.slide_height
    tbl = slide.shapes.add_table(rows, cols, Cm(0), Cm(0), slide_w, slide_h).table

    # proportional column widths
    col_props = [3, 6, 2, 3, 3]
    total = sum(col_props)
    slide_width_cm = slide_w / Cm(1)
    for i, prop in enumerate(col_props):
        tbl.columns[i].width = Cm(slide_width_cm * prop / total)

    # header row = 1 cm; data rows share remaining height
    tbl.rows[0].height = Cm(1.0)
    per_row_cm = (slide_h / Cm(1)) / rows
    for r_idx in range(1, rows):
        tbl.rows[r_idx].height = Cm(per_row_cm)

    for r_idx, row in enumerate(data):
        for c_idx, txt in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x4F, 0x81, 0xBD) if r_idx == 0 else RGBColor(0xFF, 0xFF, 0xFF)

            tf = cell.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = str(txt)
            if r_idx == 0:
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.RIGHT if c_idx in (2, 3, 4) else PP_ALIGN.LEFT

            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(12 if r_idx == 0 else 10)
            run.font.bold = (r_idx == 0)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r_idx == 0 else RGBColor(0x00, 0x00, 0x00)


def modify_existing_odp(input_path: str,
                        table_data: list,
                        summary_json: dict,
                        output_path: str):
    prs = Presentation(input_path)
    add_data_slide1(prs, summary_json)
    modify_slide2(prs, table_data)
    prs.save(output_path)






